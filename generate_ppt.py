"""
generate_ppt.py — Generate PursuitIQ Enterprise Architecture Presentation

Creates a professional .pptx with:
  - Title slide
  - Problem & Solution slide
  - 6 Intelligence Outputs slide
  - Full Architecture Diagram (single slide, fan-out agents)
  - Agent Detail slide
  - OpenAI Platform Features slide
  - Tech Stack & Cost slide
  - Closing slide

Run: python generate_ppt.py
Output: PursuitIQ_Architecture.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR

# ─── Color Palette (Dark theme, purple accent) ───────────────────────────────
BG_DARK = RGBColor(0x0F, 0x0F, 0x1A)       # Near-black background
BG_CARD = RGBColor(0x1A, 0x1A, 0x2E)       # Card background
PURPLE = RGBColor(0x7C, 0x3A, 0xED)        # Primary purple
PURPLE_LIGHT = RGBColor(0xA7, 0x8B, 0xFA)  # Light purple
GREEN = RGBColor(0x10, 0xB9, 0x81)         # Success green
RED = RGBColor(0xEF, 0x44, 0x44)           # Alert red
AMBER = RGBColor(0xF5, 0x9E, 0x0B)        # Warning amber
WHITE = RGBColor(0xFF, 0xFF, 0xFF)         # White text
GRAY = RGBColor(0x9C, 0xA3, 0xAF)         # Secondary text
BLUE = RGBColor(0x38, 0xBD, 0xF8)         # Info blue

SLIDE_WIDTH = Inches(13.333)  # Widescreen 16:9
SLIDE_HEIGHT = Inches(7.5)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=12,
                 bold=False, color=WHITE, alignment=PP_ALIGN.LEFT,
                 font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_rounded_rect(slide, left, top, width, height, fill_color, text="",
                     font_size=10, font_color=WHITE, bold=False, border_color=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()

    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_top = Pt(4)
        tf.margin_bottom = Pt(4)
        tf.margin_left = Pt(6)
        tf.margin_right = Pt(6)
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = bold
        p.font.name = "Segoe UI"
        p.alignment = PP_ALIGN.CENTER
    return shape


def add_arrow_down(slide, cx, top, length, color=PURPLE_LIGHT):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.DOWN_ARROW, cx - Inches(0.15), top, Inches(0.3), length
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def create_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Use blank layout for all slides
    blank_layout = prs.slide_layouts[6]

    # ═══════════════════════════════════════════════════════════════════════════
    # SLIDE 1: Title
    # ═══════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_DARK)

    add_text_box(slide, Inches(1), Inches(1.8), Inches(11), Inches(1.2),
                 "PursuitIQ", font_size=54, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1), Inches(3.0), Inches(11), Inches(0.8),
                 "Pursuit Intelligence Platform", font_size=28, bold=False,
                 color=PURPLE_LIGHT, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(2), Inches(4.2), Inches(9), Inches(1.5),
                 "11 Autonomous AI Agents  •  Fan-Out Parallel Execution  •  Under 8 Minutes\n"
                 "Real-Time Web Intelligence  •  GPT-5.5  •  7 OpenAI Platform Features",
                 font_size=16, color=GRAY, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(3), Inches(6.2), Inches(7), Inches(0.5),
                 "HCLTech Hackathon 2026  •  Enterprise AI Architecture",
                 font_size=13, color=GRAY, alignment=PP_ALIGN.CENTER)

    # ═══════════════════════════════════════════════════════════════════════════
    # SLIDE 2: Problem → Solution
    # ═══════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_DARK)

    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
                 "What If You Had 11 AI Researchers on Every Deal?",
                 font_size=30, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    # Problem side
    add_rounded_rect(slide, Inches(0.5), Inches(1.3), Inches(5.8), Inches(5.5),
                     BG_CARD, border_color=RED)
    add_text_box(slide, Inches(0.8), Inches(1.5), Inches(5.2), Inches(0.5),
                 "THE PROBLEM", font_size=14, bold=True, color=RED,
                 alignment=PP_ALIGN.CENTER)
    problem_text = (
        "• Pursuit teams spend 3–5 DAYS on research\n"
        "  before writing a single line\n\n"
        "• Client intel — manual web trawling\n\n"
        "• Competitor tracking — outdated databases\n\n"
        "• Pricing — last quarter's spreadsheets\n\n"
        "• RFP decoding — skim-reading 100+ pages\n\n"
        "• Hidden disqualifiers — MISSED entirely"
    )
    add_text_box(slide, Inches(1.0), Inches(2.2), Inches(5.0), Inches(4.5),
                 problem_text, font_size=14, color=GRAY)

    # Solution side
    add_rounded_rect(slide, Inches(7.0), Inches(1.3), Inches(5.8), Inches(5.5),
                     BG_CARD, border_color=GREEN)
    add_text_box(slide, Inches(7.3), Inches(1.5), Inches(5.2), Inches(0.5),
                 "PURSUITIQ", font_size=14, bold=True, color=GREEN,
                 alignment=PP_ALIGN.CENTER)
    solution_text = (
        "• 11 AI agents running in PARALLEL\n\n"
        "• Live web search — not stale data\n\n"
        "• Real Azure + AWS pricing APIs\n\n"
        "• Ghost bids — simulate competitor proposals\n\n"
        "• Hidden disqualifiers found on page 47\n\n"
        "• Under 8 minutes  •  Under $6 per run\n\n"
        "• Your team gets a 3-day head start"
    )
    add_text_box(slide, Inches(7.3), Inches(2.2), Inches(5.2), Inches(4.5),
                 solution_text, font_size=14, color=WHITE)

    # ═══════════════════════════════════════════════════════════════════════════
    # SLIDE 3: 6 Intelligence Outputs
    # ═══════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_DARK)

    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
                 "Under 8 Minutes — Your Team Gets This",
                 font_size=28, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    outputs = [
        ("RFP Decoded", "Every requirement, hidden\ndisqualifier, evaluation weight,\ndeadline — nothing missed", PURPLE),
        ("Client Signals", "Live web search — leadership\nchanges, hiring patterns,\nearnings, strategy moves", BLUE),
        ("Competitor Playbook", "Who's bidding, strengths,\npredicted positioning,\nsimulated ghost bids", RED),
        ("Real-Time Pricing", "Live Azure + AWS rates\nfrom APIs — 3 options\nwith cost breakdowns", GREEN),
        ("Win Strategy", "Killer differentiators,\nwin themes, how to beat\neach competitor", AMBER),
        ("Proposal Draft", "All intelligence structured\ninto a draft — writers\nstart with substance", PURPLE_LIGHT),
    ]

    for i, (title, desc, accent) in enumerate(outputs):
        col = i % 3
        row = i // 3
        left = Inches(0.6 + col * 4.2)
        top = Inches(1.3 + row * 3.0)

        add_rounded_rect(slide, left, top, Inches(3.8), Inches(2.6), BG_CARD,
                         border_color=accent)
        add_text_box(slide, left + Inches(0.2), top + Inches(0.2),
                     Inches(3.4), Inches(0.5), title,
                     font_size=15, bold=True, color=accent, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, left + Inches(0.2), top + Inches(0.8),
                     Inches(3.4), Inches(1.6), desc,
                     font_size=12, color=GRAY, alignment=PP_ALIGN.CENTER)

    # ═══════════════════════════════════════════════════════════════════════════
    # SLIDE 4: ARCHITECTURE DIAGRAM (Single Slide — Full Fan-Out)
    # ═══════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_DARK)

    add_text_box(slide, Inches(0.3), Inches(0.1), Inches(12.7), Inches(0.5),
                 "11-Agent Architecture — Fan-Out Parallel Execution Pipeline",
                 font_size=20, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    # ── Phase 1: RFP Input + Agent 1 ──
    add_rounded_rect(slide, Inches(5.5), Inches(0.65), Inches(2.4), Inches(0.45),
                     PURPLE, "UPLOAD RFP (PDF/DOCX)", font_size=9, bold=True)

    add_arrow_down(slide, Inches(6.7), Inches(1.12), Inches(0.22))

    # Agent 1
    add_rounded_rect(slide, Inches(4.8), Inches(1.37), Inches(3.8), Inches(0.55),
                     BG_CARD, "AGENT 1: RFP DECOMPOSER\nGPT-5.5 | Structured Outputs | reasoning: MEDIUM",
                     font_size=8, border_color=PURPLE, bold=False)

    add_arrow_down(slide, Inches(6.7), Inches(1.95), Inches(0.22))

    # ── Phase 1.5: Planner + Pre-fetch ──
    # Planner
    add_rounded_rect(slide, Inches(4.0), Inches(2.2), Inches(2.8), Inches(0.5),
                     BG_CARD, "PLANNER AGENT\nAutonomous Strategy Decisions",
                     font_size=8, border_color=PURPLE_LIGHT, bold=False)

    # Pre-fetch
    add_rounded_rect(slide, Inches(7.2), Inches(2.2), Inches(2.8), Inches(0.5),
                     BG_CARD, "PRE-FETCH\nSEC EDGAR | UK Contracts | US SAM.gov",
                     font_size=8, border_color=GRAY, bold=False)

    # Phase label
    add_text_box(slide, Inches(0.2), Inches(2.25), Inches(1.5), Inches(0.4),
                 "PHASE 1.5\n[30-70s]", font_size=7, color=GRAY, bold=True)
    add_text_box(slide, Inches(0.2), Inches(1.42), Inches(1.5), Inches(0.4),
                 "PHASE 1\n[0-30s]", font_size=7, color=GRAY, bold=True)

    # Arrow down from planner
    add_arrow_down(slide, Inches(6.7), Inches(2.73), Inches(0.22))

    # ── Phase 2: 5-Agent Fan-Out ──
    add_text_box(slide, Inches(0.2), Inches(3.0), Inches(1.5), Inches(0.4),
                 "PHASE 2\n[70-160s]", font_size=7, color=GRAY, bold=True)

    # "FAN-OUT" label
    add_text_box(slide, Inches(1.4), Inches(2.97), Inches(2.0), Inches(0.3),
                 "═══ PARALLEL FAN-OUT ═══", font_size=8, color=GREEN, bold=True)

    phase2_agents = [
        ("AGENT 2\nWin Intelligence\nFile Search (RAG)\nVector Store", GREEN),
        ("AGENT 3\nClient Intel\nWeb Search\n(Live Web)", BLUE),
        ("AGENT 4\nCompetitor\nWar Room\nWeb + SEC + Jobs", RED),
        ("DEAL\nFINGERPRINT\nPattern Match\nProcurement DB", AMBER),
        ("JOB INTEL\nLinkedIn/Indeed\nStaffing Signals\n(Enrichment)", GRAY),
    ]

    fan_start = Inches(1.5)
    fan_width = Inches(2.1)
    fan_gap = Inches(0.15)

    for i, (text, accent) in enumerate(phase2_agents):
        left = fan_start + i * (fan_width + fan_gap)
        add_rounded_rect(slide, left, Inches(3.3), fan_width, Inches(1.0),
                         BG_CARD, text, font_size=7.5, border_color=accent)

    # ── Convergence arrows ──
    add_arrow_down(slide, Inches(6.7), Inches(4.35), Inches(0.22))

    # ── Phase 3: 4-Agent Fan-Out ──
    add_text_box(slide, Inches(0.2), Inches(4.6), Inches(1.5), Inches(0.4),
                 "PHASE 3\n[160-250s]", font_size=7, color=GRAY, bold=True)

    add_text_box(slide, Inches(1.4), Inches(4.58), Inches(2.0), Inches(0.3),
                 "═══ PARALLEL FAN-OUT ═══", font_size=8, color=GREEN, bold=True)

    phase3_agents = [
        ("AGENT 5\nSolution + Pricing\nAzure/AWS Live APIs\nreasoning: HIGH", GREEN),
        ("AGENT 6\nProposal Draft\nCodex CLI (GPT-5)\nZero API Cost", PURPLE_LIGHT),
        ("GHOST BID\nSimulation\nRed Team Analysis\nWrite Their Proposal", RED),
        ("QUALITY GATE\nAutonomous Evaluator\nAccept / Retry / Deepen\nTriggers Re-runs", AMBER),
    ]

    fan3_start = Inches(1.8)
    fan3_width = Inches(2.5)
    fan3_gap = Inches(0.2)

    for i, (text, accent) in enumerate(phase3_agents):
        left = fan3_start + i * (fan3_width + fan3_gap)
        add_rounded_rect(slide, left, Inches(4.9), fan3_width, Inches(1.0),
                         BG_CARD, text, font_size=7.5, border_color=accent)

    # ── Post Phase 3: Reflection ──
    add_arrow_down(slide, Inches(6.7), Inches(5.95), Inches(0.2))

    add_text_box(slide, Inches(0.2), Inches(6.15), Inches(1.5), Inches(0.4),
                 "POST-3\nConditional", font_size=7, color=GRAY, bold=True)

    add_rounded_rect(slide, Inches(4.3), Inches(6.18), Inches(4.7), Inches(0.48),
                     BG_CARD, "REFLECTION LOOP — Self-Correction (if pricing sanity fails) → Re-runs Agent 5 autonomously",
                     font_size=8, border_color=AMBER)

    # ── Final output ──
    add_arrow_down(slide, Inches(6.7), Inches(6.7), Inches(0.2))

    add_rounded_rect(slide, Inches(3.5), Inches(6.95), Inches(6.3), Inches(0.45),
                     PURPLE, "PURSUIT INTELLIGENCE → Your Team Takes Over (Auto-Learn → Knowledge Base)",
                     font_size=9, bold=True)

    # ── OpenAI tech strip at bottom ──
    tech_strip = "Structured Outputs  |  Web Search  |  File Search  |  Vector Stores  |  Reasoning Control  |  Multi-Model  |  Codex CLI"
    add_text_box(slide, Inches(1.5), Inches(7.1), Inches(10.5), Inches(0.3),
                 tech_strip, font_size=8, color=PURPLE_LIGHT, alignment=PP_ALIGN.CENTER)

    # ═══════════════════════════════════════════════════════════════════════════
    # SLIDE 5: Agent Details Table
    # ═══════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_DARK)

    add_text_box(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
                 "11 Agents — Models, Tools, and Execution",
                 font_size=24, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    agents_data = [
        ("1", "RFP Decomposer", "GPT-5.5", "MEDIUM", "Structured Outputs (parse)", "Phase 1"),
        ("2", "Planner Agent", "GPT-5.5", "MEDIUM", "Structured Outputs (parse)", "Phase 1.5"),
        ("3", "Win Intelligence", "GPT-5.5", "MEDIUM", "File Search, Azure AI Search, TED", "Phase 2"),
        ("4", "Client Intelligence", "GPT-5.5", "MEDIUM", "Web Search (Responses API)", "Phase 2"),
        ("5", "Competitor War Room", "GPT-5.5", "MEDIUM", "Web Search, SEC EDGAR, Jobs", "Phase 2"),
        ("6", "Deal Fingerprint", "GPT-5.5", "MEDIUM", "Procurement DBs, Knowledge Base", "Phase 2"),
        ("7", "Solution + Pricing", "GPT-5.5", "HIGH", "Azure API, AWS API, KB, Procurement", "Phase 3"),
        ("8", "Proposal Draft", "GPT-5 (Codex)", "MEDIUM", "Codex CLI, KB Style Templates", "Phase 3"),
        ("9", "Ghost Bid Simulation", "GPT-5.5", "MEDIUM", "Financial data, Job intel", "Phase 3"),
        ("10", "Quality Gate", "GPT-5.5", "MEDIUM", "Evaluator (no tools)", "Phase 3"),
        ("11", "Reflection Loop", "GPT-5.5", "HIGH", "Self-assessment (conditional)", "Post-3"),
    ]

    # Headers
    headers = ["#", "Agent", "Model", "Reasoning", "Tools / APIs", "Phase"]
    col_widths = [Inches(0.4), Inches(2.0), Inches(1.3), Inches(1.0), Inches(3.5), Inches(1.2)]
    col_starts = [Inches(0.8)]
    for w in col_widths[:-1]:
        col_starts.append(col_starts[-1] + w)

    header_top = Inches(0.9)
    for i, (header, left, width) in enumerate(zip(headers, col_starts, col_widths)):
        add_text_box(slide, left, header_top, width, Inches(0.35),
                     header, font_size=10, bold=True, color=PURPLE_LIGHT)

    # Separator line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0.8), Inches(1.25), Inches(9.4), Pt(1.5))
    line.fill.solid()
    line.fill.fore_color.rgb = PURPLE
    line.line.fill.background()

    # Data rows
    row_top = Inches(1.35)
    for row_idx, row_data in enumerate(agents_data):
        y = row_top + row_idx * Inches(0.52)
        row_color = WHITE if row_idx % 2 == 0 else GRAY
        for col_idx, (cell, left, width) in enumerate(zip(row_data, col_starts, col_widths)):
            add_text_box(slide, left, y, width, Inches(0.45),
                         cell, font_size=9, color=row_color)

    # ═══════════════════════════════════════════════════════════════════════════
    # SLIDE 6: OpenAI Platform Features
    # ═══════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_DARK)

    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
                 "7 OpenAI Platform Features in Production",
                 font_size=26, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    features = [
        ("1. Structured Outputs", "response_format=Pydantic — enforces typed JSON.\nZero free text between agents. Schema-validated contracts.", PURPLE),
        ("2. Web Search", "Responses API with web_search tool.\nLive client signals, competitor news, earnings calls.", BLUE),
        ("3. File Search (RAG)", "Responses API with file_search tool.\nRAG over 100+ historical deals in Vector Store.", GREEN),
        ("4. Vector Stores", "Persistent embeddings of deal corpus.\nSimilarity matching for win patterns.", AMBER),
        ("5. Reasoning Control", "reasoning_effort = low / medium / high.\nRight compute per task complexity.", PURPLE_LIGHT),
        ("6. Multi-Model Routing", "GPT-5.5 (agents) + GPT-5 (Codex) + GPT-4.1-mini (light).\nOptimal cost/capability per agent.", RED),
        ("7. Codex CLI", "GPT-5 via ChatGPT subscription — zero API cost.\nAgent 6 proposal drafting (primary path).", GREEN),
    ]

    for i, (title, desc, accent) in enumerate(features):
        col = i % 2
        row = i // 2
        left = Inches(0.5 + col * 6.5)
        top = Inches(1.2 + row * 1.55)

        add_rounded_rect(slide, left, top, Inches(6.0), Inches(1.35), BG_CARD,
                         border_color=accent)
        add_text_box(slide, left + Inches(0.2), top + Inches(0.1),
                     Inches(5.6), Inches(0.35), title,
                     font_size=13, bold=True, color=accent)
        add_text_box(slide, left + Inches(0.2), top + Inches(0.5),
                     Inches(5.6), Inches(0.7), desc,
                     font_size=10, color=GRAY)

    # ═══════════════════════════════════════════════════════════════════════════
    # SLIDE 7: Anti-Hallucination + Resilience
    # ═══════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_DARK)

    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
                 "Anti-Hallucination Stack & Autonomous Behaviors",
                 font_size=26, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    # Anti-hallucination layers
    add_text_box(slide, Inches(0.5), Inches(1.1), Inches(6.0), Inches(0.4),
                 "5-Layer Anti-Hallucination", font_size=15, bold=True, color=RED)

    layers = [
        ("Layer 1:", "Structured Outputs — agents cannot return free-form data"),
        ("Layer 2:", "Per-agent anti-hallucination rules in system prompts"),
        ("Layer 3:", "Quality Gate — evaluates and can trigger re-runs"),
        ("Layer 4:", "Reflection Loop — pricing self-correction on sanity failure"),
        ("Layer 5:", "Verifier — cross-checks all claims against data sources"),
    ]
    for i, (label, desc) in enumerate(layers):
        y = Inches(1.6 + i * 0.5)
        add_text_box(slide, Inches(0.7), y, Inches(1.2), Inches(0.4),
                     label, font_size=11, bold=True, color=AMBER)
        add_text_box(slide, Inches(1.9), y, Inches(5.0), Inches(0.4),
                     desc, font_size=11, color=WHITE)

    # Autonomous behaviors
    add_text_box(slide, Inches(7.0), Inches(1.1), Inches(6.0), Inches(0.4),
                 "Agentic Autonomous Behaviors", font_size=15, bold=True, color=GREEN)

    behaviors = [
        "Planner autonomously identifies competitors (not hardcoded)",
        "Client Intel falls back to RFP inference if web search empty",
        "Reflection detects magnitude errors → re-runs pricing",
        "Quality Gate issues accept / retry / deepen verdicts",
        "Orchestrator times out slow agents, proceeds with available data",
        "Auto-Learn: every pursuit feeds back into knowledge base",
        "Win Intel caps confidence when evidence is sparse",
        "Ghost Bid writes competitor proposals before they do",
    ]
    for i, text in enumerate(behaviors):
        y = Inches(1.6 + i * 0.48)
        add_text_box(slide, Inches(7.2), y, Inches(5.5), Inches(0.4),
                     f"• {text}", font_size=10, color=WHITE)

    # Resilience box
    add_rounded_rect(slide, Inches(0.5), Inches(5.3), Inches(12.3), Inches(1.8),
                     BG_CARD, border_color=BLUE)
    add_text_box(slide, Inches(0.8), Inches(5.4), Inches(5.0), Inches(0.4),
                 "Fault Tolerance", font_size=14, bold=True, color=BLUE)
    resilience_text = (
        "• OpenAI client: max_retries=2 + Tenacity (3 attempts, exponential backoff 2s→30s)\n"
        "• Retries on: 429 Rate Limit, API Timeout, Connection Error\n"
        "• Per-phase timeout: 300s — slow agents skipped, pipeline continues\n"
        "• Critical agents (2, 3): failure halts pipeline  •  Non-critical (Fingerprint, Ghost Bid): failure logged, skipped\n"
        "• Codex CLI unavailable → automatic fallback to OpenAI API (GPT-5.5)"
    )
    add_text_box(slide, Inches(0.8), Inches(5.8), Inches(11.8), Inches(1.2),
                 resilience_text, font_size=10, color=GRAY)

    # ═══════════════════════════════════════════════════════════════════════════
    # SLIDE 8: Tech Stack + Cost
    # ═══════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_DARK)

    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
                 "Technology Stack & Economics",
                 font_size=26, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    # Tech stack
    add_text_box(slide, Inches(0.5), Inches(1.0), Inches(6.0), Inches(0.4),
                 "Technology Stack", font_size=15, bold=True, color=PURPLE_LIGHT)

    stack_items = [
        ("Frontend:", "Next.js 15, React 19, Tailwind CSS"),
        ("API:", "FastAPI (Python), Uvicorn ASGI"),
        ("AI Engine:", "OpenAI GPT-5.5 + GPT-5 (Codex CLI)"),
        ("Orchestration:", "ThreadPoolExecutor (6 workers, fan-out)"),
        ("Schemas:", "Pydantic v2 + Structured Outputs"),
        ("RAG:", "OpenAI Vector Stores + Azure AI Search"),
        ("Knowledge:", "Azure Blob Storage + AI Search"),
        ("Cloud Pricing:", "Azure Retail API + AWS Pricing API (live)"),
        ("Procurement:", "TED Europe, Contracts Finder, SAM.gov, AusTender, GeBIZ"),
        ("Financial:", "SEC EDGAR (margins, revenue, pricing floors)"),
    ]
    for i, (label, value) in enumerate(stack_items):
        y = Inches(1.5 + i * 0.42)
        add_text_box(slide, Inches(0.7), y, Inches(1.8), Inches(0.35),
                     label, font_size=10, bold=True, color=AMBER)
        add_text_box(slide, Inches(2.5), y, Inches(4.5), Inches(0.35),
                     value, font_size=10, color=WHITE)

    # Cost breakdown
    add_text_box(slide, Inches(7.0), Inches(1.0), Inches(6.0), Inches(0.4),
                 "Cost Per Pursuit Run", font_size=15, bold=True, color=GREEN)

    costs = [
        ("GPT-5.5 (10 agents)", "~$5.00"),
        ("Codex CLI / GPT-5 (Agent 6)", "$0.00"),
        ("Azure/AWS Pricing APIs", "$0.00"),
        ("OpenAI Web Search", "~$0.30"),
        ("File Search (Vector Store)", "~$0.10"),
        ("Procurement APIs", "$0.00"),
        ("", ""),
        ("TOTAL PER PURSUIT", "~$5.50"),
    ]
    for i, (item, cost) in enumerate(costs):
        y = Inches(1.5 + i * 0.42)
        bold = (i == len(costs) - 1)
        color = GREEN if bold else (WHITE if item else WHITE)
        add_text_box(slide, Inches(7.2), y, Inches(3.5), Inches(0.35),
                     item, font_size=11 if bold else 10, bold=bold, color=color)
        add_text_box(slide, Inches(10.7), y, Inches(1.5), Inches(0.35),
                     cost, font_size=11 if bold else 10, bold=bold, color=color,
                     alignment=PP_ALIGN.RIGHT)

    # ROI Box
    add_rounded_rect(slide, Inches(7.0), Inches(5.0), Inches(5.5), Inches(2.0),
                     BG_CARD, border_color=GREEN)
    add_text_box(slide, Inches(7.3), Inches(5.1), Inches(5.0), Inches(0.4),
                 "ROI", font_size=14, bold=True, color=GREEN)
    roi_text = (
        "Manual research: 3–5 days × multiple people\n"
        "PursuitIQ: Under 8 minutes × $5.50\n\n"
        "That's a 500x speed improvement\n"
        "for less than the cost of a coffee."
    )
    add_text_box(slide, Inches(7.3), Inches(5.5), Inches(5.0), Inches(1.5),
                 roi_text, font_size=12, color=WHITE)

    # ═══════════════════════════════════════════════════════════════════════════
    # SLIDE 9: Closing
    # ═══════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_DARK)

    add_text_box(slide, Inches(1), Inches(2.0), Inches(11), Inches(1.0),
                 "11 Agents. 7 OpenAI Features.\nFully Autonomous.",
                 font_size=40, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(2), Inches(3.8), Inches(9), Inches(0.8),
                 "Under 8 minutes  •  Under $6  •  3-day head start",
                 font_size=22, color=PURPLE_LIGHT, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(2), Inches(5.0), Inches(9), Inches(1.5),
                 "Your team still decides. Still strategizes. Still writes.\n"
                 "They just start with days of research already done.",
                 font_size=16, color=GRAY, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(3), Inches(6.5), Inches(7), Inches(0.5),
                 "PursuitIQ — Pursuit Intelligence Platform",
                 font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)

    # ═══════════════════════════════════════════════════════════════════════════
    # SAVE
    # ═══════════════════════════════════════════════════════════════════════════
    output_path = "PursuitIQ_Architecture.pptx"
    prs.save(output_path)
    print(f"\n[OK] Presentation saved: {output_path}")
    print(f"  Slides: {len(prs.slides)}")
    print(f"  Format: 16:9 Widescreen")
    print(f"  Theme: Dark (purple accent)")
    return output_path


if __name__ == "__main__":
    create_presentation()
