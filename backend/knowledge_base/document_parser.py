"""
document_parser.py — Extract structured text from proposal documents.

Supports PDF, DOCX, and PPTX. Extracts by section/page for intelligent chunking.
"""

import io
import logging
from typing import Optional

logger = logging.getLogger(__name__)


def extract_text(file_bytes: bytes, filename: str) -> list[dict]:
    """
    Extract text from a document, split into sections.

    Returns list of chunks:
        [{"section": "Executive Summary", "content": "...", "page": 3}, ...]
    """
    ext = filename.lower().rsplit(".", 1)[-1] if "." in filename else ""

    if ext == "pdf":
        return _extract_pdf(file_bytes)
    elif ext == "docx":
        return _extract_docx(file_bytes)
    elif ext in ("pptx", "ppt"):
        return _extract_pptx(file_bytes)
    else:
        logger.warning(f"Unsupported file type: {ext}, treating as plain text")
        return [{"section": "Full Document", "content": file_bytes.decode("utf-8", errors="ignore"), "page": 1}]


def _extract_pdf(file_bytes: bytes) -> list[dict]:
    """Extract text from PDF with page markers, grouped by detected sections."""
    import fitz

    doc = fitz.open(stream=file_bytes, filetype="pdf")
    chunks = []
    current_section = "Cover / Introduction"
    current_content = []
    current_page = 1

    section_keywords = [
        "executive summary", "table of contents", "scope", "approach",
        "methodology", "solution", "architecture", "pricing", "commercial",
        "team", "credentials", "references", "why choose", "differentiator",
        "delivery", "timeline", "risk", "governance", "appendix",
        "technical approach", "management approach", "staffing",
    ]

    for page_num, page in enumerate(doc, start=1):
        text = page.get_text("text").strip()
        if not text:
            continue

        lines = text.split("\n")
        for line in lines:
            line_lower = line.strip().lower()
            is_heading = (
                len(line.strip()) < 80 and
                any(kw in line_lower for kw in section_keywords) and
                (line.strip().isupper() or line.strip().istitle() or line.strip()[0:2].isdigit())
            )

            if is_heading and current_content:
                chunks.append({
                    "section": current_section,
                    "content": "\n".join(current_content).strip(),
                    "page": current_page,
                })
                current_section = line.strip()
                current_content = []
                current_page = page_num
            else:
                current_content.append(line)

    if current_content:
        chunks.append({
            "section": current_section,
            "content": "\n".join(current_content).strip(),
            "page": current_page,
        })

    doc.close()

    if not chunks:
        all_text = ""
        doc2 = fitz.open(stream=file_bytes, filetype="pdf")
        for page in doc2:
            all_text += page.get_text("text") + "\n"
        doc2.close()
        chunks = [{"section": "Full Document", "content": all_text.strip(), "page": 1}]

    return chunks


def _extract_docx(file_bytes: bytes) -> list[dict]:
    """Extract text from DOCX, grouped by headings."""
    from docx import Document

    doc = Document(io.BytesIO(file_bytes))
    chunks = []
    current_section = "Introduction"
    current_content = []

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue

        if para.style.name.startswith("Heading"):
            if current_content:
                chunks.append({
                    "section": current_section,
                    "content": "\n".join(current_content).strip(),
                    "page": None,
                })
            current_section = text
            current_content = []
        else:
            current_content.append(text)

    if current_content:
        chunks.append({
            "section": current_section,
            "content": "\n".join(current_content).strip(),
            "page": None,
        })

    # Also extract from tables
    for table in doc.tables:
        table_text = []
        for row in table.rows:
            row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if row_text:
                table_text.append(" | ".join(row_text))
        if table_text:
            chunks.append({
                "section": "Table Data",
                "content": "\n".join(table_text),
                "page": None,
            })

    return chunks if chunks else [{"section": "Full Document", "content": "\n".join(p.text for p in doc.paragraphs), "page": 1}]


def _extract_pptx(file_bytes: bytes) -> list[dict]:
    """Extract text from PowerPoint presentations."""
    from pptx import Presentation

    prs = Presentation(io.BytesIO(file_bytes))
    chunks = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_text = []
        slide_title = f"Slide {slide_num}"

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        if shape == slide.shapes.title or (hasattr(shape, 'placeholder_format') and shape.placeholder_format and shape.placeholder_format.idx == 0):
                            slide_title = text
                        else:
                            slide_text.append(text)

            if shape.has_table:
                for row in shape.table.rows:
                    row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if row_text:
                        slide_text.append(" | ".join(row_text))

        if slide_text:
            chunks.append({
                "section": slide_title,
                "content": "\n".join(slide_text),
                "page": slide_num,
            })

    return chunks if chunks else [{"section": "Empty Presentation", "content": "", "page": 1}]


def get_full_text(file_bytes: bytes, filename: str) -> str:
    """Get the complete text of a document as a single string."""
    chunks = extract_text(file_bytes, filename)
    return "\n\n".join(
        f"[{c['section']}]\n{c['content']}"
        for c in chunks
        if c['content'].strip()
    )
